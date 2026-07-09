# -*- coding: utf-8 -*-
"""약손 최종 발표(4분) .pptx 생성 — 편집 가능한 네이티브 텍스트."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x0f, 0x1e, 0x3d)
GREEN = RGBColor(0x2f, 0x6a, 0x4f)
INK   = RGBColor(0x1f, 0x29, 0x37)
GRAY  = RGBColor(0x55, 0x65, 0x70)
WHITE = RGBColor(0xff, 0xff, 0xff)
RED   = RGBColor(0xb9, 0x1c, 0x1c)
FONT  = "맑은 고딕"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def band(s, num, title):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.15))
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background(); r.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.2), SW - Inches(1.2), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = f"{num}.  {title}"
    run.font.size = Pt(30); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = FONT


def body(s, top=1.55):
    tb = s.shapes.add_textbox(Inches(0.7), Inches(top), SW - Inches(1.4), SH - Inches(top + 0.5))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def fill(tf, lines):
    # lines: (text, size, bold, color, bullet, space_after)
    for i, (text, size, bold, color, bullet, space) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        run = p.add_run()
        run.text = ("•  " + text) if bullet else text
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = FONT


# ===== 슬라이드 1 — 표지 + 한 줄 소개 =====
s = slide()
tb = s.shapes.add_textbox(Inches(1), Inches(2.1), SW - Inches(2), Inches(3.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "💊 약손"
r.font.size = Pt(72); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = FONT
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(12)
r = p.add_run(); r.text = "약속하는 약손"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = FONT
p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(22)
r = p.add_run(); r.text = '"눈이 어두운 어르신의 약 복용 문제를,\n찍으면 읽어주는 약손으로 해결합니다."'
r.font.size = Pt(20); r.font.color.rgb = GRAY; r.font.name = FONT
notes(s, "[1] 표지 + 한 줄 소개(Hook)  ⏱️0:20\n\n"
         "안녕하세요. 저희 앱은 약손입니다. "
         "눈이 어두워 약봉투를 못 읽는 어르신의 복약 문제를, 찍으면 읽어주는 약손으로 해결합니다.")

# ===== 슬라이드 2 — 문제 & 왜 지금 =====
s = slide(); band(s, 2, "문제 & '왜 지금'")
fill(body(s), [
    ("● 무엇이 문제인가", 22, True, GREEN, False, 6),
    ("약봉투 작은 글씨 → 무슨 약인지, 하루 몇 번 먹는지 헷갈림", 20, False, INK, True, 6),
    ("잘못 복용·거름 → 복약사고, 건강 악화", 20, False, INK, True, 6),
    ("혼자 사는 어르신은 물어볼 사람도 없음", 20, False, INK, True, 16),
    ("● 왜 지금인가 (초고령사회)", 22, True, GREEN, False, 6),
    ("2024년 말 대한민국 초고령사회 진입 — 65세 이상 20% 돌파  (※수치 확인)", 20, False, INK, True, 6),
    ("노인 대부분이 만성질환으로 여러 약을 동시에 복용(다제약물)  (※수치 확인)", 20, False, INK, True, 6),
])
notes(s, "[2] 문제 & 왜 지금  ⏱️0:40\n\n"
         "약봉투 글씨는 너무 작습니다. 무슨 약인지, 하루 몇 번 먹는지 헷갈려 잘못 드시거나 거르게 되고, "
         "이는 복약사고로 이어집니다. 혼자 사시는 분은 물어볼 사람도 없죠. "
         "그리고 지금, 우리나라는 이미 초고령사회에 진입했고 어르신 대부분이 여러 약을 한꺼번에 드십니다. "
         "이 문제는 앞으로 더 커집니다.\n\n※ 2번 통계는 발표 전 실제 최신 수치로 교체하세요.")

# ===== 슬라이드 3 — 타깃 페르소나 =====
s = slide(); band(s, 3, "타깃 — 페르소나")
fill(body(s), [
    ("👵 이복순 · 74세 · 1인 가구  (주 사용자)", 22, True, INK, False, 6),
    ("전화·카톡은 하지만 작은 글씨는 안 보임", 19, False, GRAY, True, 4),
    ("정형외과·내과에서 한 번에 3~4개 약 처방 → 매번 헷갈림", 19, False, GRAY, True, 16),
    ("👩 김미영 · 45세 · 이복순 님의 딸  (지불자)", 22, True, INK, False, 6),
    ("타지 거주, 부모님 복약이 늘 걱정", 19, False, GRAY, True, 4),
    ('"약 잘 드시는지"만 알 수 있다면 기꺼이 돈 낼 의향', 19, False, GRAY, True, 16),
    ("핵심: 쓰는 사람(어르신)과 돈 내는 사람(자녀)이 다르다", 21, True, GREEN, False, 0),
])
notes(s, "[3] 타깃 페르소나  ⏱️0:35\n\n"
         "핵심 사용자는 혼자 사시는 74세 이복순 님입니다. 여러 약을 받아오시지만 매번 헷갈리시죠. "
         "그리고 멀리 사는 딸 김미영 씨는 늘 불안합니다. "
         "여기서 중요한 건, 쓰는 사람과 돈 내는 사람이 다르다는 점입니다.")

# ===== 슬라이드 4 — 솔루션 & 핵심기능 (I-P-O) =====
s = slide(); band(s, 4, "솔루션 & 핵심 기능  (입력 → 처리 → 출력)")
# 세 박스
box_w, box_h, top = Inches(3.9), Inches(2.0), Inches(1.7)
labels = [
    ("입력 (Input)", ["약봉투를", "카메라로 찍기 한 번"], NAVY),
    ("처리 (Process)", ["글자 인식 →", "시니어 눈높이로 정리", "(이름·효능·복용법·하루 몇 번)"], GREEN),
    ("출력 (Output)", ["큰 글씨 + 음성 + 쉬운 설명"], RGBColor(0xc2,0x41,0x0c)),
]
xs = [Inches(0.7), Inches(4.72), Inches(8.74)]
for (title, lines, col), x in zip(labels, xs):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    for ln in lines:
        pp = tf.add_paragraph(); pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = ln; rr.font.size = Pt(15); rr.font.color.rgb = WHITE; rr.font.name = FONT
# 출력 상세
tf2 = body(s, top=3.95)
fill(tf2, [
    ("출력으로 더해지는 돌봄 기능", 20, True, GREEN, False, 6),
    ("⏰ 복약 알림(아침·낮·저녁)  ·  💊 약 소진 3일 전 알림", 19, False, INK, True, 5),
    ("👨‍👩‍👧 가족 공유(최대 5명)  ·  🏥 주변 병원 찾기  ·  🆘 긴급 SOS", 19, False, INK, True, 0),
])
notes(s, "[4] 솔루션 & 핵심기능 (I-P-O)  ⏱️0:45\n\n"
         "동작은 입력-처리-출력 세 단계입니다. 입력은 약봉투를 찍는 것, 딱 한 번. "
         "처리는 글자를 읽어 시니어 눈높이의 쉬운 말로 정리하고, 출력으로 큰 글씨와 음성으로 알려드립니다. "
         "여기에 복약 알림, 약이 떨어지기 3일 전 알림, 가족 공유, 주변 병원 찾기, 긴급 SOS까지 담았습니다.")

# ===== 슬라이드 5 — 기존 서비스와 무엇이 다른가 =====
s = slide(); band(s, 5, "기존 서비스와 무엇이 다른가")
rows, cols = 4, 3
tbl = s.shapes.add_table(rows, cols, Inches(0.7), Inches(1.7), SW - Inches(1.4), Inches(3.6)).table
data = [
    ["구분", "기존 (검색·약학정보)", "약손"],
    ["사용법", "약 이름 타이핑 검색", "찍기 한 번"],
    ["대상", "일반인", "시니어 전용 (큰 글씨·음성·쉬운 말)"],
    ["범위", "약 하나씩 검색", "처방전 통째로 + 알림·가족·병원·SOS"],
]
for ci, w in enumerate([Inches(2.0), Inches(4.6), Inches(5.13)]):
    tbl.columns[ci].width = w
for ri in range(rows):
    for cj in range(cols):
        cell = tbl.cell(ri, cj)
        cell.text = data[ri][cj]
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]; run.font.name = FONT
        if ri == 0:
            run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            run.font.size = Pt(17); run.font.color.rgb = INK
            run.font.bold = (cj == 2)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xff,0xff,0xff) if cj != 2 else RGBColor(0xe9,0xf7,0xef)
notes(s, "[5] 기존 서비스와 차이  ⏱️0:30\n\n"
         "기존 서비스는 약 이름을 직접 검색해야 하고 일반인 대상입니다. "
         "약손은 찍기 한 번, 시니어 전용, 그리고 단순 정보가 아니라 알림·가족 돌봄·긴급 SOS까지 아우르는 통합 돌봄이라는 점이 다릅니다.")

# ===== 슬라이드 6 — 프로토타입 시연 =====
s = slide(); band(s, 6, "프로토타입 시연 (실제 작동하는 웹 앱)")
fill(body(s), [
    ("약.jpg로 라이브 데모 — 시연 순서", 22, True, GREEN, False, 10),
    ("① 약봉투 찍기 → 약 4개를 큰 글씨 + 음성으로, 하루 몇 번까지", 20, False, INK, True, 7),
    ("② 복약 알림 → 시계가 뜨고 진동·소리로 알려줌", 20, False, INK, True, 7),
    ("③ 가족 관리 → 등록한 가족의 약을 함께 확인", 20, False, INK, True, 7),
    ("④ 주변 병원 → GPS로 가까운 병원 거리순 + 길찾기", 20, False, INK, True, 7),
    ("⑤ SOS → 5초 꾹 → 사이렌 + 현재 위치와 함께 119", 20, False, INK, True, 0),
])
notes(s, "[6] 프로토타입 시연  ⏱️0:50  ★핵심\n\n"
         "실제 작동하는 화면을 보여드리겠습니다. (찍고) 약 네 개를 큰 글씨와 음성으로, 하루 몇 번 먹는지까지 알려줍니다. "
         "복약 시간엔 시계가 울리고, 가족의 약도 함께 챙기고, 내 위치 주변 병원을 거리순으로 찾아 길안내하고, "
         "위급하면 SOS를 5초 누르면 위치와 함께 119로 연결됩니다.\n\n"
         "※ 미리 약.jpg 준비 + 리허설. 소리·GPS는 localhost에서 시연.")

# ===== 슬라이드 7 — 기대효과 & 마무리 =====
s = slide(); band(s, 7, "기대 효과 & 마무리")
fill(body(s), [
    ("기대 효과", 22, True, GREEN, False, 6),
    ("어르신 — 혼자서도 정확한 복약 → 복약사고↓, 자립↑", 20, False, INK, True, 5),
    ("자녀 — 멀리서도 안심", 20, False, INK, True, 5),
    ("사회 — 복약순응도↑ → 의료비 절감", 20, False, INK, True, 16),
    ("향후 계획", 22, True, GREEN, False, 6),
    ("실제 의약품 DB 연동 · 약국/지자체 제휴 · 푸시 알림", 20, False, INK, True, 18),
    ("약을 챙기겠다는 약속 — 약속하는 약손", 24, True, NAVY, False, 0),
])
notes(s, "[7] 기대효과 & 마무리  ⏱️0:20\n\n"
         "약손은 어르신께 자립을, 자녀에게 안심을, 사회에는 의료비 절감을 드립니다. "
         "앞으로 의약품 DB 연동과 약국·지자체 제휴로 키워가겠습니다. "
         "약을 챙기겠다는 약속, 약속하는 약손이었습니다. 감사합니다.")

out = r"C:\Users\USER\Desktop\프로젝트\약손_발표.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
